"""Load and process AxleWave company documents."""
from pathlib import Path
from typing import List, Dict
from zipfile import ZipFile
import xml.etree.ElementTree as ET


def extract_docx_text(docx_path: Path) -> str:
    """Extract text from .docx file."""
    with ZipFile(docx_path) as docx:
        xml_content = docx.read('word/document.xml')
        tree = ET.XML(xml_content)
        
        namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
        paragraphs = []
        
        for paragraph in tree.findall('.//w:p', namespaces):
            texts = [node.text for node in paragraph.findall('.//w:t', namespaces) if node.text]
            if texts:
                paragraphs.append(''.join(texts))
        
        return '\n'.join(paragraphs)


def extract_pdf_text(pdf_path: Path) -> str:
    """Extract text from PDF file."""
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(str(pdf_path))
        text = []
        for page in reader.pages:
            text.append(page.extract_text())
        return '\n'.join(text)
    except Exception as e:
        print(f"PDF extraction error: {e}")
        return ""


def extract_xlsx_text(xlsx_path: Path) -> str:
    """Extract text from Excel file."""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(xlsx_path), data_only=True)
        text = []
        for sheet in wb.worksheets:
            text.append(f"Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    text.append(row_text)
        return '\n'.join(text)
    except Exception as e:
        print(f"Excel extraction error: {e}")
        return ""


def extract_pptx_text(pptx_path: Path) -> str:
    """Extract text from PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    except Exception as e:
        print(f"PowerPoint extraction error: {e}")
        return ""


def load_axlewave_documents(docs_dir: Path) -> List[Dict[str, str]]:
    """Load all AxleWave documents from directory."""
    documents = []
    
    extractors = {
        '.docx': extract_docx_text,
        '.pdf': extract_pdf_text,
        '.xlsx': extract_xlsx_text,
        '.pptx': extract_pptx_text
    }
    
    for ext, extractor in extractors.items():
        for file_path in docs_dir.glob(f"*{ext}"):
            try:
                content = extractor(file_path)
                if content.strip():
                    documents.append({
                        "filename": file_path.name,
                        "content": content,
                        "type": ext[1:]
                    })
            except Exception as e:
                print(f"Error loading {file_path.name}: {e}")
    
    return documents


def create_company_context(documents: List[Dict[str, str]]) -> str:
    """Create a consolidated context about AxleWave from all documents."""
    context_parts = []
    
    for doc in documents:
        context_parts.append(f"=== {doc['filename']} ===\n{doc['content']}\n")
    
    return "\n".join(context_parts)


class DocumentLoader:
    """Simple wrapper for document loading functions."""
    
    @staticmethod
    def load_documents(docs_dir: Path) -> List[Dict[str, str]]:
        return load_axlewave_documents(docs_dir)
    
    @staticmethod
    def create_context(documents: List[Dict[str, str]]) -> str:
        return create_company_context(documents)
